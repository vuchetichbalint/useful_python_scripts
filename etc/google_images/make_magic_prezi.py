import pandas as pd

from pptx import Presentation
from pptx.util import Inches
from PIL import Image

from googlescraper import get_google_images

def get_ratio(w,h,w0,h0):
	print(w,h,w0,h0)
	# it resizes to maximum 640*360pxs
	if w < w0 and h < h0:
		print('kicsi volt!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!')
		return w,h
	if w < w0:
		print('kicsi volt!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!')
		return w*h0//h, h0
	if h < h0:
		print('kicsi volt!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!')
		return w0, h*w0//w
	if (w/w0 < h/h0):
		print('elso if ter vissza')
		return w*h0//h, h0
	else:
		print('else-nel lepett ki')
		return w0, h*w0//w

def resize_rational(file):
	w0 = 640
	h0 = 360
	
	im = Image.open(file)
	if im.size[0] < w0 and im.size[1] < h0:
		return
	
	im = im.resize(get_ratio(im.size[0], im.size[1], w0, h0), Image.ANTIALIAS)
	# see docs: http://effbot.org/imagingbook/image.htm#tag-Image.Image.size
	im.save(file)

def get_filename(term):
	return ''.join(filter(str.isalpha, term))

def generate_prezi(df):
	prs = Presentation()

	for idx, row in df.iterrows():
		title_only_slide_layout = prs.slide_layouts[5]
		slide = prs.slides.add_slide(title_only_slide_layout)

		shapes = slide.shapes
		title_shape = shapes.title

		title_shape.text = row['title']

		img_path = 'img/' + row['filename']
		left = top = Inches(2)
		pic = slide.shapes.add_picture(img_path, left, top)

		pic.left = (prs.slide_width - pic.width) // 2

	prs.save('prezi/generated_prezi.pptx')

if __name__ == "__main__":
	df = pd.read_csv('terms.csv', sep=';')
	# ha nincs mas akkor a termet adja cimnek
	df = df.fillna(method='ffill', axis='columns')
	df['filename'] = df['term'].apply(lambda x: get_google_images(x, 4, get_filename(x))[-1])
	df['filename'].apply(lambda x: resize_rational('img/' + x))
	generate_prezi(df)




"""
be kell olvasni a szavakat egy file-bol
masodik keresni kell hozza kepeket (ezt a v2-ben csinalhatnank random is)
letrehozni egy ppt-t es belepakolni a kepeket a szovegekkel egyutt

input format:
szavak:

[kifejezes,cim]
ha nincs cim, akkor a kifejezest teszi be cimnek
"""


