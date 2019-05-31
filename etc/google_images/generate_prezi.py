'''

from pptx import Presentation
from pptx.util import Inches


def generate_prezi(df):
	prs = Presentation()

	for idx, row in df.iterrows():
		title_only_slide_layout = prs.slide_layouts[5]
		slide = prs.slides.add_slide(title_only_slide_layout)

		shapes = slide.shapes
		title_shape = shapes.title

		title_shape.text = row['title']

		img_path = 'img/' + row['filename']
		left = top = Inches(2)
		pic = slide.shapes.add_picture(img_path, left, top)

		pic.left = (prs.slide_width - pic.width) // 2

	prs.save('prezi/generated_prezi.pptx')

'''